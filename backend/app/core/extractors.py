from __future__ import annotations
from io import BytesIO
from typing import Tuple, Callable, Optional

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

TEXT_MIME_PREFIXES = ("text/",)
TEXT_EXTS = (
    ".txt", ".md", ".csv", ".json", ".log", ".py", ".js", ".ts", ".tsx", ".jsx",
    ".html", ".css", ".sql", ".yaml", ".yml"
)


PDF_MIME = "application/pdf"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

IMAGE_MIME_PREFIXES = ("image/",)
IMAGE_EXTS = (
    ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"
)

# Optional OCR dependencies (only used if installed)
try:
    from PIL import Image  # type: ignore
except Exception:  # pragma: no cover
    Image = None  # type: ignore

try:
    import pytesseract  # type: ignore
except Exception:  # pragma: no cover
    pytesseract = None  # type: ignore


def _ocr_image_bytes(data: bytes) -> str:
    """OCR an image byte payload to text.

    Requires Pillow + pytesseract to be installed. If not available, returns "".
    """
    if Image is None or pytesseract is None:
        return ""

    try:
        img = Image.open(BytesIO(data))
        # normalize to RGB for OCR
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        return (pytesseract.image_to_string(img) or "").strip()
    except Exception:
        return ""


def _safe_decode(raw: bytes) -> str:
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("latin-1", errors="replace")
    
OCRFN = Callable[[bytes], str]

def extract_text_from_pdf(data:bytes) -> Tuple[str, str]:
    reader =  PdfReader(BytesIO(data))
    parts: list[str] = []

    for page in reader.pages:
        t = page.extract_text() or ""
        t = t.strip()
        if t:
            parts.append(t)

    text = "\n\n".join(parts).strip()
    if not text:
        return "needs_ocr", ""
    return "extracted" , text

def extract_text_from_docx(data: bytes) -> Tuple[str, str]:
    doc = Document(BytesIO(data))
    parts: list[str] = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            parts.append(t)

    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join((cell.text or "").strip() for cell in row.cells).strip()
            if row_text:
                parts.append(row_text)

    text = "\n".join(parts).strip()
    return ("extracted", text) if text else ("extracted", "")   


def extract_text_from_pptx(data: bytes) -> Tuple[str, str]:
    prs = Presentation(BytesIO(data))
    parts: list[str] = []

    for si, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            # shape.text exists only for some shapes
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    slide_parts.append(t)
        if slide_parts:
            parts.append(f"[Slide {si}]\n" + "\n".join(slide_parts))

    text = "\n\n".join(parts).strip()
    return ("extracted", text) if text else ("extracted", "")


def extract_text_from_textlike(data: bytes) -> Tuple[str, str]:
    text = _safe_decode(data).strip()
    return ("extracted", text) if text else ("extracted", "")

def extract_text_from_image(data: bytes) -> Tuple[str, str]:
    # Try built-in OCR if Pillow + pytesseract are available.
    text = _ocr_image_bytes(data)
    if text:
        return ("ocr_extracted", text)
    return ("needs_ocr", "")


def pick_extractor(filename: str, mime: str):
    name = (filename or "").lower()

    # PDF
    if mime == PDF_MIME or name.endswith(".pdf"):
        return extract_text_from_pdf

    # DOCX
    if mime == DOCX_MIME or name.endswith(".docx"):
        return extract_text_from_docx

    # PPTX
    if mime == PPTX_MIME or name.endswith(".pptx"):
        return extract_text_from_pptx

    # Text-like (by mime prefix or extension)
    if mime.startswith(TEXT_MIME_PREFIXES) or name.endswith(TEXT_EXTS):
        return extract_text_from_textlike
    
    if mime.startswith(IMAGE_MIME_PREFIXES) or name.endswith(IMAGE_EXTS):
        return extract_text_from_image

    # Unknown/binary: accept, but no extraction
    return None


def extract_text_any(
        *, filename: str, mime: str, data: bytes, ocr: Optional[OCRFN] = None
) -> Tuple[str, str]:
    extractor = pick_extractor(filename, mime)
    if extractor is None:
        return "unknown_format", ""
    
    status, text = extractor(data)
    if status == "needs_ocr":
        # If caller provided a custom OCR function, prefer it.
        if ocr is not None:
            ocr_text = (ocr(data) or "").strip()
            return ("ocr_extracted", ocr_text) if ocr_text else ("ocr_failed", "")

        # Otherwise, for images we may have already tried built-in OCR; for PDFs/scans we still need OCR.
        return ("needs_ocr", "")
    
    return status, text